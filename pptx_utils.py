from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_THEME_COLOR

from utils import hex_to_rgb # 假設 utils.py 在同一目錄下

def change_slide_background(presentation: Presentation, hex_color: str):
    """
    變更簡報中所有投影片的背景顏色。

    Args:
        presentation (Presentation): python-pptx Presentation 物件。
        hex_color (str): 目標背景顏色的 Hex 字串 (例如 "#RRGGBB")。
    """
    rgb_color = hex_to_rgb(hex_color)
    if rgb_color == (0, 0, 0) and hex_color.strip().lower() not in ["#000000", "#000", "#000000ff"]: # 檢查是否是預設的黑色，避免誤判
        print(f"警告：'{hex_color}' 無法轉換為有效的 RGB 顏色，將使用預設顏色。")

    for slide in presentation.slides:
        background = slide.background
        fill = background.fill
        fill.solid() # 設定為純色填充
        fill.fore_color.rgb = RGBColor(*rgb_color)

def change_text_style(presentation: Presentation, title_color_hex: str, body_color_hex: str):
    """
    遍歷簡報中的所有形狀，並根據其類型（標題或內文）更改文字顏色。
    """
    title_rgb = hex_to_rgb(title_color_hex)
    body_rgb = hex_to_rgb(body_color_hex)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = False
            # --- 修正開始：使用 if shape.is_placeholder 進行安全檢查 ---
            if shape.is_placeholder:
                # 現在可以安全地存取 placeholder_format
                if shape.placeholder_format.type in [
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                    PP_PLACEHOLDER.SUBTITLE
                ]:
                    is_title = True
            # --- 修正結束 ---
            
            target_rgb = title_rgb if is_title else body_rgb

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*target_rgb)

def _get_slide_content(slide):
    """
    安全地從投影片中提取標題和主要內文。此版本更穩健，包含後備邏輯。
    """
    title = ""
    content = ""
    title_shape = None

    if slide.shapes.title:
        title = slide.shapes.title.text
        title_shape = slide.shapes.title

    content_placeholders = [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]
    
    max_text_len_placeholder = -1
    for shape in slide.shapes:
        # --- 修正開始：使用 if shape.is_placeholder 進行安全檢查 ---
        if shape.is_placeholder and shape.has_text_frame:
            if shape.placeholder_format.type in content_placeholders:
                if len(shape.text) > max_text_len_placeholder:
                    content = shape.text
                    max_text_len_placeholder = len(shape.text)
        # --- 修正結束 ---

    if not content:
        max_text_len_fallback = -1
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            
            if shape.has_text_frame:
                if len(shape.text) > max_text_len_fallback:
                    max_text_len_fallback = len(shape.text)
                    content = shape.text

    if not title and content:
        title = content.split('\n')[0]

    return {'title': title, 'content': content}

def apply_layout_to_slides(presentation: Presentation, layout_index: int):
    """
    使用新的版面配置重新建立投影片，並將它們附加到簡報的末尾。
    """
    if not (0 <= layout_index < len(presentation.slide_layouts)):
        print(f"錯誤：無效的版面索引 {layout_index}")
        return

    target_layout = presentation.slide_layouts[layout_index]
    
    original_slides_content = []
    slides_to_process = list(presentation.slides)
    for slide in slides_to_process:
        content = _get_slide_content(slide)
        if content['title'] or content['content']:
            original_slides_content.append(content)

    if not original_slides_content:
        print("警告：未能從原始簡報中提取到任何有效內容。")
        return

    try:
        separator_layout = next(l for l in presentation.slide_layouts if "Header" in l.name or "Title Only" in l.name)
        sep_slide = presentation.slides.add_slide(separator_layout)
        sep_slide.shapes.title.text = "--- AI 重新設計的版面 ---"
    except StopIteration:
        print("找不到合適的分隔頁版面，將直接新增投影片。")

    for item in original_slides_content:
        new_slide = presentation.slides.add_slide(target_layout)
        
        if new_slide.shapes.title:
            new_slide.shapes.title.text = item['title']
        
        body_shape = None
        for shape in new_slide.placeholders:
            # --- 修正開始：使用 if shape.is_placeholder 進行安全檢查 ---
            if shape.is_placeholder:
                if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    body_shape = shape
                    break
            # --- 修正結束 ---
        
        if body_shape:
            body_shape.text = item['content']
        else:
            print(f"警告：目標版面 '{target_layout.name}' 沒有找到合適的內文預留位置。")
